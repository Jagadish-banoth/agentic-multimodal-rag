"""
ingestion/text_ingest.py
Day 3: Production-quality text ingestion + token-aware chunking

Usage:
  - Put PDFs/TXT files in the raw_dir configured in config/settings.yaml
  - Run: python ingestion/text_ingest.py
Outputs:
  - data/processed/chunks.jsonl   (one JSON chunk per line)
  - data/processed/chunks_index.json (manifest & stats)
"""

import os
import re
import json
import yaml
import hashlib
import logging
from pathlib import Path
from typing import List, Tuple, Dict
from dataclasses import dataclass, asdict
from tqdm import tqdm
import fitz  # PyMuPDF

try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    from docx import Document
except Exception:
    Document = None
try:
    import openpyxl
except Exception:
    openpyxl = None
try:
    from bs4 import BeautifulSoup
except Exception:
    BeautifulSoup = None
try:
    from PIL import Image
except Exception:
    Image = None
try:
    import pytesseract
except Exception:
    pytesseract = None
import io
import csv
from transformers import AutoTokenizer
import regex as re2  # more powerful regex

# ---------- CONFIG ----------
DEFAULT_CONFIG = {
    "ingestion": {
        "raw_dir": "data/raw/text",
        "out_dir": "data/processed",
        "chunk_tokens": 512,
        "overlap_tokens": 92,
        "sentence_trimming": True,
        "remove_headers_footers": True,
    },
    "models": {
        # prefer 'tokenizer_model' key; keep compatibility with old 'tokenizer_name'
        "tokenizer_model": "nomic-ai/nomic-embed-text-v1.5"
    },
    "logging": {"level": "INFO", "logfile": "logs/text_ingest.log"},
}


# ---------- UTIL TYPES ----------
@dataclass
class ChunkObj:
    chunk_id: str
    content: str
    source: str
    page_start: int
    page_end: int
    modality: str = "text"
    token_count: int = 0
    char_start: int = 0
    char_end: int = 0
    meta: Dict = None


# ---------- SETUP ----------
def load_config(path: str = "config/settings.yaml"):
    cfg = DEFAULT_CONFIG.copy()
    if Path(path).exists():
        with open(path, "r", encoding="utf-8") as f:
            user_cfg = yaml.safe_load(f)
            # shallow merge simple keys
            for k, v in user_cfg.items():
                if isinstance(v, dict):
                    cfg.get(k, {}).update(v)
                else:
                    cfg[k] = v
    return cfg


cfg = load_config()
RAW_DIR = Path(cfg["ingestion"]["raw_dir"])
OUT_DIR = Path(cfg["ingestion"]["out_dir"])
OUT_DIR.mkdir(parents=True, exist_ok=True)
CHUNKS_OUT = OUT_DIR / "chunks.jsonl"
INDEX_OUT = OUT_DIR / "chunks_index.json"

LOGFILE = cfg["logging"].get("logfile", "logs/text_ingest.log")
os.makedirs(Path(LOGFILE).parent, exist_ok=True)

logging.basicConfig(
    filename=LOGFILE,
    filemode="a",
    level=getattr(logging, cfg["logging"].get("level", "INFO")),
    format="%(asctime)s - %(levelname)s - %(message)s",
)
console = logging.StreamHandler()
console.setLevel(getattr(logging, cfg["logging"].get("level", "INFO")))
logging.getLogger().addHandler(console)
logger = logging.getLogger(__name__)


# ---------- Helpers ----------
def list_input_files(raw_dir: Path) -> List[Path]:
    exts = [
        "*.pdf",
        "*.txt",
        "*.pptx",
        "*.docx",
        "*.md",
        "*.html",
        "*.csv",
        "*.xlsx",
        "*.json",  # Add JSON support
        "*.jpg",
        "*.jpeg",
        "*.png",
        "*.tif",
        "*.tiff",
    ]
    files = []
    for e in exts:
        files.extend(sorted(raw_dir.glob(e)))
    return files


def normalize_whitespace(text: str) -> str:
    # remove soft hyphenations and join
    text = text.replace("\u00ad", "")  # soft hyphen
    text = re.sub(r"-\n", "", text)  # word-split across lines
    # normalize newlines
    text = re.sub(r"\r\n?", "\n", text)
    # collapse many blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # trim trailing spaces
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def extract_text_from_pdf(path: Path) -> List[Tuple[int, str]]:
    doc = fitz.open(str(path))
    pages = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        text = page.get_text("text")
        # if no text, try OCR rendering to image then pytesseract
        if not text or not text.strip():
            if pytesseract and Image:
                try:
                    pix = page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_bytes))
                    ocr_text = pytesseract.image_to_string(img)
                    text = ocr_text
                except Exception:
                    text = ""
        pages.append((i + 1, text))
    return pages


def extract_text_from_txt(path: Path) -> List[Tuple[int, str]]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [(1, text)]


def extract_text_from_docx(path: Path) -> List[Tuple[int, str]]:
    if Document is None:
        raise RuntimeError(
            "python-docx is not installed. Install with 'pip install python-docx'"
        )
    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text]
    return [(1, "\n".join(parts))]


def extract_text_from_xlsx(path: Path) -> List[Tuple[int, str]]:
    if openpyxl is None:
        raise RuntimeError(
            "openpyxl is not installed. Install with 'pip install openpyxl'"
        )
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(" ".join([str(c) for c in row if c is not None]))
        parts.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
    return [(1, "\n\n".join(parts))]


def extract_text_from_html(path: Path) -> List[Tuple[int, str]]:
    if BeautifulSoup is None:
        raise RuntimeError(
            "beautifulsoup4 is not installed. Install with 'pip install beautifulsoup4'"
        )
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n")
    return [(1, text)]


def extract_text_from_csv(path: Path) -> List[Tuple[int, str]]:
    text_lines = []
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            text_lines.append(" ".join(row))
    return [(1, "\n".join(text_lines))]

def extract_text_from_json(path: Path) -> List[Tuple[int, str]]:
    """
    Extracts text from a JSON file.
    Expects a list of objects, where each object has a "text", "content", or "page_content" field.
    """
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    
    pages = []
    if isinstance(data, list):
        for idx, item in enumerate(data):
            text = item.get("text") or item.get("content") or item.get("page_content")
            if text:
                pages.append((idx + 1, text))
    return pages


# --- JSON QA extraction ---
def extract_qa_from_json(path: Path) -> List[Tuple[int, str]]:
    """
    Expects a JSON file containing a list of objects with 'question' and 'answer' fields.
    Returns a list of (index, text) tuples, where text = 'Q: ...\nA: ...'.
    """
    with path.open("r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    pages = []
    if isinstance(data, dict) and "data" in data:
        data = data["data"]
    for idx, item in enumerate(data):
        q = item.get("question") or item.get("Q")
        a = item.get("answer") or item.get("A")
        if q and a:
            text = f"Q: {q}\nA: {a}"
            pages.append((idx + 1, text))
    return pages


def extract_text_from_image(path: Path) -> List[Tuple[int, str]]:
    if pytesseract is None or Image is None:
        raise RuntimeError(
            "pytesseract and pillow are required for image OCR. Install with 'pip install pytesseract pillow' and install Tesseract on your system."
        )
    img = Image.open(str(path))
    text = pytesseract.image_to_string(img)
    return [(1, text)]


def extract_text_from_pptx(path: Path) -> List[Tuple[int, str]]:
    """Extract text from PPTX slides; each slide becomes a page tuple (slide_index, text)."""
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is not installed. Install with 'pip install python-pptx'"
        )
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                try:
                    txt = shape.text
                except Exception:
                    txt = ""
                if txt:
                    parts.append(txt)
        slides.append((i + 1, "\n".join(parts)))
    return slides


def detect_and_remove_headers_footers(page_texts: List[str]) -> List[str]:
    """
    Heuristic: find common first-line and last-line across pages using similarity,
    remove them when they appear frequently.
    """
    from difflib import SequenceMatcher

    def most_common_similar(lines):
        # compute similarity matrix and find candidate common lines
        lines = [s.strip() for s in lines if s and s.strip()]
        if not lines:
            return None
        freq = {}
        for a in lines:
            for b in lines:
                if a == b:
                    freq[a] = freq.get(a, 0) + 1
                else:
                    sim = SequenceMatcher(None, a[:60], b[:60]).ratio()
                    if sim > 0.95:
                        freq[a] = freq.get(a, 0) + 1
        if not freq:
            return None
        candidate, count = max(freq.items(), key=lambda kv: kv[1])
        # set threshold relative to pages
        if count >= max(2, 0.5 * len(lines)):
            return candidate
        return None

    first_lines = []
    last_lines = []
    for p in page_texts:
        s = p.strip().split("\n")
        if s:
            first_lines.append(s[0])
            last_lines.append(s[-1])
    header = most_common_similar(first_lines)
    footer = most_common_similar(last_lines)
    cleaned = []
    for p in page_texts:
        lines = p.split("\n")
        if (
            header
            and lines
            and lines[0].strip()
            and (
                lines[0].strip().startswith(header[:10])
                or header.startswith(lines[0].strip()[:10])
            )
        ):
            lines = lines[1:]
        if (
            footer
            and lines
            and lines[-1].strip()
            and (
                lines[-1].strip().endswith(footer[-10:])
                or footer.endswith(lines[-1].strip()[-10:])
            )
        ):
            lines = lines[:-1]
        cleaned.append("\n".join(lines).strip())
    return cleaned


# Sentence splitter: try nltk if available, else regex fallback
try:
    import nltk

    nltk.data.find("tokenizers/punkt")

    def split_sentences(text):
        return nltk.sent_tokenize(text)

except Exception:
    # simple regex-based sentence splitter (works reasonably well)
    SENTENCE_END_RE = re2.compile(r'(?<=\S[.!?])\s+(?=[A-Z0-9"\'“])')

    def split_sentences(text):
        # split on sentence-ending punctuation followed by whitespace and capital letter/number/quote
        if not text:
            return []
        parts = SENTENCE_END_RE.split(text)
        return [p.strip() for p in parts if p.strip()]


# deterministic chunk id
def deterministic_chunk_id(source: str, char_start: int, char_end: int) -> str:
    key = f"{source}|{char_start}|{char_end}"
    h = hashlib.sha1(key.encode("utf-8")).hexdigest()[:16]
    safe = Path(source).name.replace(" ", "_")
    return f"{safe}_{char_start}_{char_end}_{h}"


# ---------- Tokenizer loader ----------
def load_tokenizer(name: str):
    try:
        tokenizer = AutoTokenizer.from_pretrained(name, use_fast=True, trust_remote_code=True)
        logger.info(f"Loaded tokenizer: {name}")
        return tokenizer
    except Exception as e:
        logger.exception(f"Failed loading tokenizer {name}: {e}")
        raise


# ---------- Chunking core ----------
def chunk_document(
    full_text: str,
    page_offsets: List[Tuple[int, int, int]],
    tokenizer,
    chunk_size: int,
    overlap: int,
    sentence_trimming: bool = True,
) -> List[ChunkObj]:
    """
    full_text: concatenation of pages with separators; page_offsets = list of (page_num, char_start, char_end)
    Returns list of ChunkObj
    """
    enc = tokenizer(full_text, add_special_tokens=False, return_offsets_mapping=True)
    if "offset_mapping" not in enc or not enc["offset_mapping"]:
        logger.warning(
            "Tokenizer offset_mapping not available; falling back to whitespace tokenization."
        )
        # fallback: simple whitespace tokens
        tokens = full_text.split()
        token_count = len(tokens)
        step = chunk_size - overlap
        chunks = []
        for s in range(0, token_count, step):
            e = min(s + chunk_size, token_count)
            chunk_text = " ".join(tokens[s:e])
            # approximate char positions by searching (best-effort)
            cs = full_text.find(chunk_text)
            ce = cs + len(chunk_text) if cs >= 0 else 0
            chunk_id = deterministic_chunk_id("fallback", cs, ce)
            chunks.append(
                ChunkObj(
                    chunk_id,
                    chunk_text,
                    "inline",
                    1,
                    1,
                    "text",
                    e - s,
                    cs,
                    ce,
                    {"fallback": True},
                )
            )
            if e == token_count:
                break
        return chunks

    token_offsets = enc["offset_mapping"]
    token_count = len(token_offsets)
    step = max(1, chunk_size - overlap)
    chunks_out: List[ChunkObj] = []

    for start_token in range(0, token_count, step):
        end_token = min(start_token + chunk_size, token_count)
        # compute char span
        start_char = token_offsets[start_token][0] if token_offsets[start_token] else 0
        end_char = (
            token_offsets[end_token - 1][1]
            if token_offsets[end_token - 1]
            else len(full_text)
        )
        if start_char >= end_char:
            continue
        chunk_text = full_text[start_char:end_char].strip()
        # sentence trimming: ensure chunk ends at sentence boundary when enabled
        if sentence_trimming:
            sents = split_sentences(chunk_text)
            if len(sents) >= 2:
                # If last sentence seems partial (very short) drop it
                last = sents[-1]
                if len(last.split()) < 6:
                    sents = sents[:-1]
                    candidate = " ".join(sents).strip()
                    if candidate:
                        chunk_text = candidate
                        # adjust end_char to reflect new length
                        # naive: find last occurrence of chunk_text ending at or before end_char
                        idx = full_text.rfind(
                            chunk_text, max(0, start_char - 10), end_char
                        )
                        if idx >= 0:
                            start_char = idx
                            end_char = idx + len(chunk_text)
        # map page span
        page_start = page_end = None
        for pnum, pstart, pend in page_offsets:
            # overlap test
            if start_char >= pstart and start_char < pend:
                page_start = pnum
            if end_char > pstart and end_char <= pend:
                page_end = pnum
        if page_start is None:
            page_start = page_offsets[0][0] if page_offsets else 1
        if page_end is None:
            page_end = page_offsets[-1][0] if page_offsets else page_start

        token_window_count = end_token - start_token
        chunk_id = deterministic_chunk_id("doc", start_char, end_char)
        c = ChunkObj(
            chunk_id=chunk_id,
            content=chunk_text,
            source="",  # fill at caller
            page_start=page_start,
            page_end=page_end,
            modality="text",
            token_count=token_window_count,
            char_start=start_char,
            char_end=end_char,
            meta={},
        )
        chunks_out.append(c)
        if end_token == token_count:
            break

    return chunks_out


# ---------- Document processing ----------
def process_file(
    path: Path,
    tokenizer,
    chunk_size: int,
    overlap: int,
    do_header_footer: bool,
    sentence_trimming: bool,
) -> List[ChunkObj]:
    logger.info(f"Processing {path}")
    ext = path.suffix.lower()
    if ext == ".pdf":
        pages = extract_text_from_pdf(path)
    elif ext in [".txt", ".md"]:
        pages = extract_text_from_txt(path)
    elif ext == ".pptx":
        try:
            pages = extract_text_from_pptx(path)
        except RuntimeError as e:
            logger.error(str(e))
            return []
    elif ext == ".docx":
        try:
            pages = extract_text_from_docx(path)
        except RuntimeError as e:
            logger.error(str(e))
            return []
    elif ext in [".xlsx", ".xls"]:
        try:
            pages = extract_text_from_xlsx(path)
        except RuntimeError as e:
            logger.error(str(e))
            return []
    elif ext in [".html", ".htm"]:
        try:
            pages = extract_text_from_html(path)
        except RuntimeError as e:
            logger.error(str(e))
            return []
    elif ext == ".csv":
        pages = extract_text_from_csv(path)
    elif ext == ".json":
        try:
            # Always try generic text extraction first for .json
            pages = extract_text_from_json(path)
            if not pages:
                # Fallback to QA extraction if no generic text is found
                pages = extract_qa_from_json(path)
        except Exception as e:
            logger.error(f"Failed to parse JSON file {path}: {e}")
            return []
    elif ext in [".jpg", ".jpeg", ".png", ".tif", ".tiff"]:
        try:
            pages = extract_text_from_image(path)
        except RuntimeError as e:
            logger.error(str(e))
            return []
    else:
        logger.warning(f"Unsupported file {path}")
        return []
    # normalize per page
    page_texts = [normalize_whitespace(t) for (_, t) in pages]
    if do_header_footer and len(page_texts) > 2:
        page_texts = detect_and_remove_headers_footers(page_texts)
    # build full_text and offsets
    full_text = ""
    page_offsets = []
    for pnum, ptext in zip([p[0] for p in pages], page_texts):
        start_char = len(full_text)
        # add page text and two newlines as separator
        full_text += ptext + "\n\n"
        end_char = len(full_text)
        page_offsets.append((pnum, start_char, end_char))
    if not full_text.strip():
        logger.warning(f"No text for {path}")
        return []
    raw_chunks = chunk_document(
        full_text, page_offsets, tokenizer, chunk_size, overlap, sentence_trimming
    )
    # fill source and meta
    for c in raw_chunks:
        c.source = str(path)
        c.meta = {"filename": path.name}
    logger.info(f"{path.name}: produced {len(raw_chunks)} chunks")
    return raw_chunks


# ---------- Writing ----------
def write_chunks_jsonl(chunks: List[ChunkObj], out_path: Path):
    with out_path.open("w", encoding="utf-8") as fout:
        for c in chunks:
            fout.write(json.dumps(asdict(c), ensure_ascii=False) + "\n")


def write_index(manifest: Dict, out_path: Path):
    out_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


# ---------- Run ----------
def run():
    # support both 'tokenizer_model' and legacy 'tokenizer_name'
    tokenizer_name = cfg["models"].get("tokenizer_model") or cfg["models"].get(
        "tokenizer_name"
    )
    if tokenizer_name is None:
        raise RuntimeError(
            "No tokenizer specified in config under 'models.tokenizer_model' or 'models.tokenizer_name'"
        )
    tokenizer = load_tokenizer(tokenizer_name)
    files = list_input_files(RAW_DIR)
    if not files:
        logger.error(f"No files found in {RAW_DIR}. Place PDF/TXT files and retry.")
        return
    chunk_size = cfg["ingestion"]["chunk_tokens"]
    overlap = cfg["ingestion"]["overlap_tokens"]
    sentence_trimming = cfg["ingestion"]["sentence_trimming"]
    do_header_footer = cfg["ingestion"]["remove_headers_footers"]

    all_chunks: List[ChunkObj] = []
    stats = {"files": [], "total_chunks": 0, "tokens": []}
    for f in tqdm(files, desc="Processing files"):
        chunks = process_file(
            f, tokenizer, chunk_size, overlap, do_header_footer, sentence_trimming
        )
        all_chunks.extend(chunks)
        stats["files"].append({"file": str(f), "n_chunks": len(chunks)})
        stats["total_chunks"] += len(chunks)
        stats["tokens"].extend([c.token_count for c in chunks])

    # write outputs
    logger.info(f"Writing {len(all_chunks)} chunks to {CHUNKS_OUT}")
    write_chunks_jsonl(all_chunks, CHUNKS_OUT)
    manifest = {
        "n_files": len(files),
        "n_chunks": len(all_chunks),
        "chunk_tokens": chunk_size,
        "overlap_tokens": overlap,
        "files": stats["files"],
    }
    if stats["tokens"]:
        import statistics

        manifest["token_stats"] = {
            "min": min(stats["tokens"]),
            "max": max(stats["tokens"]),
            "mean": statistics.mean(stats["tokens"]),
            "median": statistics.median(stats["tokens"]),
        }
    write_index(manifest, INDEX_OUT)
    logger.info(f"Ingestion complete. Manifest written to {INDEX_OUT}")


if __name__ == "__main__":
    run()


# --- Programmatic API for Streamlit/app integration ---
def ingest_text(uploaded_file):
    """
    Ingest a text-like file from a file-like object (e.g., Streamlit upload).
    Returns a document ID (filename) or status.
    """
    import tempfile
    from pathlib import Path

    # Save uploaded file to a temp file in RAW_DIR
    RAW_DIR.mkdir(parents=True, exist_ok=True)
    suffix = Path(uploaded_file.name).suffix or ".txt"
    with tempfile.NamedTemporaryFile(delete=False, dir=RAW_DIR, suffix=suffix) as tmp:
        tmp.write(uploaded_file.read())
        tmp_path = Path(tmp.name)

    # Use the same config as batch mode
    tokenizer_name = cfg["models"].get("tokenizer_model") or cfg["models"].get("tokenizer_name")
    tokenizer = load_tokenizer(tokenizer_name)
    chunk_size = cfg["ingestion"]["chunk_tokens"]
    overlap = cfg["ingestion"]["overlap_tokens"]
    sentence_trimming = cfg["ingestion"]["sentence_trimming"]
    do_header_footer = cfg["ingestion"]["remove_headers_footers"]

    # Process the file
    chunks = process_file(
        tmp_path, tokenizer, chunk_size, overlap, do_header_footer, sentence_trimming
    )
    if not chunks:
        logger.warning(f"No chunks produced for {uploaded_file.name}")
        return None

    # Write to processed outputs (append mode)
    write_chunks_jsonl(chunks, CHUNKS_OUT)
    # Optionally update manifest/index (not strictly needed for single file)
    # Return a document ID (filename)
    return tmp_path.name
