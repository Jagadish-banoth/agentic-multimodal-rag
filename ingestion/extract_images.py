
from pathlib import Path
from PIL import Image
import fitz  # PyMuPDF
import io
import base64
import re

from docx import Document
from pptx import Presentation

from ingestion.image_utils import save_pil_image


# -----------------------------
# PDF
# -----------------------------
def extract_from_pdf(pdf_path: Path):
    images = []
    doc = fitz.open(pdf_path)

    for page_idx, page in enumerate(doc):
        for img_idx, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base = doc.extract_image(xref)
            img_bytes = base["image"]

            pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            saved = save_pil_image(
                pil_img,
                prefix=f"{pdf_path.stem}_p{page_idx}_img{img_idx}",
            )
            images.append(saved)

    return images


# -----------------------------
# DOCX
# -----------------------------
def extract_from_docx(docx_path: Path):
    images = []
    doc = Document(docx_path)

    for idx, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.target_ref:
            img_bytes = rel.target_part.blob
            pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            saved = save_pil_image(
                pil_img,
                prefix=f"{docx_path.stem}_img{idx}",
            )
            images.append(saved)

    return images


# -----------------------------
# PPTX
# -----------------------------
def extract_from_pptx(pptx_path: Path):
    images = []
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture
                img_bytes = shape.image.blob
                pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                saved = save_pil_image(
                    pil_img,
                    prefix=f"{pptx_path.stem}_s{slide_idx}_img{shape_idx}",
                )
                images.append(saved)

    return images


# -----------------------------
# HTML (base64 embedded images)
# -----------------------------
def extract_from_html(html_path: Path):
    images = []
    text = html_path.read_text(encoding="utf-8", errors="ignore")

    matches = re.findall(
        r'<img[^>]+src="data:image/(.*?);base64,(.*?)"', text, re.DOTALL
    )

    for idx, (_, b64_data) in enumerate(matches):
        img_bytes = base64.b64decode(b64_data)
        pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        saved = save_pil_image(
            pil_img,
            prefix=f"{html_path.stem}_img{idx}",
        )
        images.append(saved)

    return images


# -----------------------------
# Direct Image Files
# -----------------------------
def extract_from_image(image_path: Path):
    pil_img = Image.open(image_path).convert("RGB")
    saved = save_pil_image(
        pil_img,
        prefix=image_path.stem,
    )
    return [saved]

def extract_images(file_path: Path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_from_pdf(file_path)
    if suffix == ".docx":
        return extract_from_docx(file_path)
    if suffix == ".pptx":
        return extract_from_pptx(file_path)
    if suffix == ".html":
        return extract_from_html(file_path)
    if suffix in [".jpg", ".jpeg", ".png", ".tif", ".tiff"]:
        return extract_from_image(file_path)

    return []